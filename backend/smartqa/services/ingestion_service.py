import os
from langchain_core.documents import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter
from smartqa.rag.vectorstore import add_documents

import pypdf
import docx
from pptx import Presentation

def process_uploaded_file(file_path, filename):
    text = ""
    print(f"Processing file: {filename} at {file_path}") # Debug Print

    try:
        ext = filename.split('.')[-1].lower()

        if ext == 'pdf':
            try:
                reader = pypdf.PdfReader(file_path)
                print(f"PDF Loaded. Pages: {len(reader.pages)}")
                for i, page in enumerate(reader.pages):
                    page_text = page.extract_text()
                    if page_text:
                        text += page_text + "\n"
                    else:
                        print(f"Warning: Page {i} empty.")
            except Exception as e:
                print(f"pypdf error: {e}")
                return {"error": f"Failed to read PDF: {e}"}
        
        elif ext in ['docx', 'doc']:
            doc = docx.Document(file_path)
            for para in doc.paragraphs:
                text += para.text + "\n"
        
        elif ext in ['pptx', 'ppt']:
            prs = Presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
        
        else:
            return {"error": f"Unsupported file type: {ext}"}

        # Check if text was found
        if not text.strip():
            print("Extraction finished but text is empty.")
            return {"error": "No text extracted. File might be scanned image or empty."}

        print(f"Extracted {len(text)} characters.")

        text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=1000,
            chunk_overlap=100
        )
        
        doc_chunks = text_splitter.create_documents(
            [text], 
            metadatas=[{"source": filename}]
        )

        add_documents(doc_chunks)

        return {"message": f"Success. Added {len(doc_chunks)} chunks from {filename}."}

    except Exception as e:
        print(f"General Processing Error: {e}")
        return {"error": f"Processing failed: {str(e)}"}
    
    finally:
        # Cleanup
        if os.path.exists(file_path):
            os.remove(file_path)